# gen_exam_focus.py
# Builds DSAI4205_Exam_Focus.docx from:
#   1. Exam Review (1).pptx  — Dr. Fong's review slides (all 146)
#   2. Transcript tips       — extracted from Apr 11 lecture

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt as PPTXPt
from docx import Document
from docx.shared import Pt, RGBColor, Cm, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

PPTX = Path('C:/Users/User/Downloads/4205/Exam Review (1).pptx')
OUT  = Path('C:/Users/User/Downloads/4205/DSAI4205_Exam_Focus.docx')

# ── Transcript-sourced tips, keyed by slide title (partial match) ────────────
TIPS = {
    'Six Vs':           'EXAM FORMAT: Given a company paragraph → identify which V\'s are present (quote evidence) → propose a solution for each V. NOT asked for bare definitions.',
    'Practice question':'MODEL ANSWER: Volume→Spark batch, Velocity→Spark Streaming, Variety→Data Lake/SparkSQL, Veracity→data validation + anomaly detection.',
    'HDFS':             'Know: 3-way replication, fixed-size partitions (128 MB), write-once-read-many. If one node fails → recover from the other two replicas.',
    'Dask':             'MCQ ONLY (1–2 questions). Focus: lazy execution, task graph, map_overlap, chunk size trade-offs. NOT in long question.',
    'Block Algorithm':  'Block algorithms split large computation into chunks. Workers operate on one chunk at a time. Fold pattern: partial result per chunk → combine.',
    'Task Graph':       'Task graph = logical plan drawn before any execution (lazy). Call .compute() to trigger. visualize() draws the DAG.',
    'Splitting Chunk':  'Small chunks = too much task scheduling overhead. Large chunks = no parallelism benefit (like not using Dask at all). Sweet spot ~100 MB–1 GB.',
    'MapReduce':        'KEY EXAM TOPIC. Know all 3 stages: Mapper → Shuffle (automatic, system-handled) → Reducer. Common mistake: students forget to mention shuffle. Intermediate format after mapper: (word, 1). After shuffle: (word, [1,1,1]). After reducer: (word, 3).',
    'Map Reduce':       'KEY EXAM TOPIC. Know all 3 stages: Mapper → Shuffle (automatic, system-handled) → Reducer. Common mistake: students forget to mention shuffle.',
    'RDD':              '10% of exam is RDD programming. Revise Tutorial 3 (Q2, Q3, Q4) and take-home inverted index exercise. Questions give a skeleton with TO-DO gaps — not blank-page coding.',
    'Resilient':        'RDD = immutable + distributed. Fault tolerance via LINEAGE (not replication). If partition lost → recompute from lineage log. Transformations are lazy; actions trigger execution.',
    'Granularity':      'Narrow transform = 1 input partition → 1 output partition (no shuffle). Wide transform = multiple input partitions → shuffle required. Know examples of each.',
    'Inverted Index':   'This take-home exercise IS a past exam question (makeup version). Pattern: flatMap (word, docID) → distinct → groupByKey → mapValues(sorted) → sortByKey.',
    'Narrow':           'Narrow: filter, map, flatMap. Wide: groupByKey, reduceByKey, join, sortByKey. Wide transforms cause network shuffle — expensive.',
    'Join and Partition':'Wide transformation example. Data must be co-partitioned for efficient join. Broadcast join avoids shuffle for small tables.',
    'DataFrame':        'DataFrame is MCQ ONLY (3–4 questions, ~4–5%). NOT in long question. Know: Catalyst optimizer, schema, key operations vs RDD.',
    'Spark SQL':        'SparkSQL = SQL interface on DataFrames. Register as temp view → query with spark.sql(). Catalyst optimiser applies predicate pushdown, column pruning.',
    'RDD vs DataFrame': 'RDD: flexible, no optimisation, good for custom logic. DataFrame: Catalyst-optimised, faster for structured data, column-level operations. Exam: DF is MCQ only.',
    'Subword Token':    'BPE exam question: target vocabulary size = 10 (not 20). Run merge rounds until vocab hits 10. Unknown character → [UNK] shown as "?" in exam.',
    'Byte Pair':        'BPE steps: (1) split words into characters. (2) count most frequent adjacent pair. (3) merge pair → new token. (4) repeat until vocab size = target (10 in exam).',
    'Tokenizing':       'To tokenise a new sentence with trained BPE vocab: greedily apply learned merges in order. If character/pair not in vocab → maps to [UNK] / "?".',
    'Stop Words':       'Stop words = high-frequency words with little meaning (the, is, at). Remove before TF-IDF or embedding. NLTK: stopwords.words("english").',
    'Stemming':         'Stemming: chop endings (Porter stemmer). Fast but may produce non-real words (studies → studi). Not guaranteed to be valid English.',
    'Lemmatization':    'Lemmatization: look up real base form (studies → study). Always produces valid English word. Slower than stemming but more accurate.',
    'One Hot':          'One-hot encoding: convert categorical string to binary vector. Problem: high dimensionality for large vocabulary. Leads to sparse representations.',
    'Bag of Words':     'BoW: count word occurrences per document, ignore order. Vocabulary = all unique words. Representation = count vector per document.',
    'Levenshtein':      'EXAM WARNING: Source word ALWAYS at TOP (columns), target word ALWAYS at LEFT (rows). "Every semester students mix this up — your whole question will be wrong." Recurrence: dp[i][j] = min(delete, insert, substitute).',
    'DP table':         'Fill Levenshtein DP left-to-right, top-to-bottom. If characters match: take diagonal (no cost). If mismatch: 1 + min(left, above, diagonal). Orientation: source=top, target=left.',
    'Word Embedding':   'NO MATH EQUATIONS in exam. Know conceptually: distributional hypothesis, static embeddings (Word2Vec) vs contextual (BERT), two embedding matrices (W and W\').',
    'Word2Vec':         'Word2Vec produces STATIC embeddings (one vector per word regardless of context). BERT/Transformers give contextual embeddings. Exam tests Word2Vec concepts.',
    'Skip-gram':        'Skip-gram: centre word → predict surrounding context words. Works better for rare words. Each training example: (centre, context) pair within window.',
    'Loss function':    'No need to memorise equations. Know WHY softmax: ensures probabilities sum to 1 and are in [0,1]. Exponent guarantees positive. Division normalises.',
    'CBOW':             'CBOW: context words → predict centre word. Average the context word vectors (divide by 2M where M = half-window). M=2 means 2 words each side, NOT 4 total context words.',
    'CBOW Example':     'CBOW forward pass: (1) look up context word embeddings. (2) average them. (3) multiply by output matrix W\'. (4) softmax → probabilities. Highest prob = predicted word.',
    'Negative Sampling':'Negative sampling avoids computing softmax over full vocabulary. For each positive (centre, context) pair, sample k random negative words. Binary logistic regression instead of full softmax.',
    'PageRank: Matrix': 'EXAM WARNING: Source node = TOP (columns). Destination node = LEFT (rows). Column sums must equal 1 (probability distribution over outlinks). Same orientation as Levenshtein.',
    'Power Iteration':  'Steps: (1) Init r = [1/N, …, 1/N]. (2) r_new = β×M×r + (1−β)/N. (3) Compute L1 change = Σ|r_new−r_old|. (4) If change < threshold → converged. Else r = r_new → repeat.',
    'dead end':         'Dead ends: pages with no outlinks leak rank (everything eventually goes to 0). Fix: treat dead-end as linking to ALL pages equally → redistributes leaked rank.',
    'Dead End':         'Dead ends: pages with no outlinks leak rank (everything eventually goes to 0). Fix: treat dead-end as linking to ALL pages equally → redistributes leaked rank.',
    'Spider Trap':      'Spider trap: self-contained cycle absorbs ALL rank. Fix: teleportation — with probability (1−β), surfer jumps to any random page regardless of links.',
    'Teleport':         'Teleportation parameter β ≈ 0.85. With prob β: follow a link. With prob (1−β): jump to random page. Fixes BOTH dead ends and spider traps simultaneously.',
    'Google Matrix':    'Google Matrix: A = βM + (1−β)(1/N)J. β ≈ 0.85. J = all-ones matrix. r = A×r at each iteration. This single formula handles both problems.',
    'undirected graph': 'Undirected graph: max edges = N(N−1)/2. Directed: max = N(N−1). Density = actual edges / max possible edges.',
    'Degree Distribut': 'P(k) = fraction of nodes with degree k. Random graph → Poisson. Scale-free (real-world: web, social) → power law P(k) ∝ k^(−γ). Most nodes few connections, few hubs have many.',
    'Paths in a graph': 'Path: sequence of nodes where each consecutive pair is connected. Simple path: no repeated nodes. Distance = length of shortest path.',
    'Distance in a Gra':'Distance = shortest path length. If no path exists (disconnected graph) → distance = ∞.',
    'Network Diameter': 'Diameter = maximum shortest path across all node pairs. Average path length = average over all pairs. "Six degrees of separation" = avg path length ≈ 6 for social networks.',
    'Clustering Coeff': 'C(v) = 2m_v / (d_v(d_v−1)). m_v = edges among v\'s neighbours. d_v = degree of v. Measures how cliquey neighbours are. High CC = tight community.',
    'Connectivity':     'Connected (undirected): path between every pair. Strongly connected (directed): path in BOTH directions. Weakly connected (directed): ignoring direction, it\'s connected.',
    'Centrality':       'Exam tests: degree, betweenness (vertex only), closeness, harmonic. NOT tested: edge betweenness ("ash centrality will not be answered in your final exam").',
    'Closeness':        'Closeness = (N−1) / Σ d(v, all others). Problem: if graph is disconnected, some distances = ∞ → formula breaks. Fix: harmonic centrality = Σ 1/d(u,v).',
    'Harmonic':         'Harmonic centrality = Σ_{u≠v} 1/d(u,v). Handles disconnected graphs (infinite distances → contribute 0, not break the formula).',
    'Betweenness':      'Betweenness(v) = Σ_{s≠v≠t} σ(s,t|v)/σ(s,t). σ(s,t) = total shortest paths from s to t. σ(s,t|v) = those passing through v. VERTEX betweenness ONLY. Edge betweenness NOT tested.',
    'Example of Betwee':'For each source s: BFS to find all shortest paths and counts. Count paths through each intermediate node. Normalise by (n−1)(n−2)/2 for undirected graphs.',
    'ACID':             'MCQ ONLY for NoSQL/CAP topics. ACID: Atomicity (all-or-nothing), Consistency, Isolation, Durability. Hard to maintain in distributed systems.',
    'NoSQL':            'NoSQL gives up some ACID guarantees in exchange for horizontal scalability and availability. Not "no SQL" — many NoSQL systems still support some query language.',
    'CAP Theorem':      'Can only guarantee 2 of 3: Consistency, Availability, Partition Tolerance. Networks always partition → choose CP (lock on partition) or AP (return stale data on partition). CA = single-node only.',
    'CA vs CP vs AP':   'CA: single-node RDBMS. CP: HBase, Zookeeper, MongoDB (strong consistency, may be unavailable during partition). AP: Cassandra, DynamoDB, CouchDB (always available, eventually consistent).',
    'BASE':             'BASE = Basically Available + Soft state + Eventually consistent. NoSQL alternative to ACID. "Soft state" = DB can be inconsistent for a period then self-heals.',
    'RDBMS vs NoSQL':   'RDBMS (CP): strong consistency, may lock and become unavailable during partition. NoSQL (AP): always responds, may return stale data, eventually becomes consistent.',
    'balance = 100':    'Strong consistency example: locks prevent reads until replication completes. Eventual consistency: reads may return old value briefly, then catches up. Instagram likes = eventual consistency is fine.',
}

# ── Section groupings (slide ranges, 1-based) ────────────────────────────────
SECTIONS = [
    ('Exam Scope & Format',                  range(1, 3)),
    ('L1 — The 6 V\'s of Big Data',          range(3, 12)),
    ('L2 — HDFS',                            range(12, 14)),
    ('L2 — Dask',                            range(14, 26)),
    ('L2/L3 — MapReduce',                    range(26, 35)),
    ('L3 — RDD & Spark Core',                range(35, 50)),
    ('L3/L4 — DataFrame & SparkSQL',         range(50, 71)),
    ('L5 — NLP: Text Preprocessing',         range(71, 92)),
    ('L6 — Word Embeddings',                 range(92, 105)),
    ('L7 — PageRank',                        range(105, 116)),
    ('L8 — Graph Analytics',                 range(116, 136)),
    ('L9 — NoSQL & CAP Theorem',             range(136, 147)),
]

# ── DOCX helpers ─────────────────────────────────────────────────────────────
def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)

def add_tip(doc, tip_text, color='warn'):
    colors = {
        'warn':    (RGBColor(0x92, 0x40, 0x0E), RGBColor(0xFF, 0xF7, 0xED)),
        'danger':  (RGBColor(0x7F, 0x1D, 0x1D), RGBColor(0xFF, 0xEE, 0xEE)),
        'info':    (RGBColor(0x1E, 0x40, 0xAF), RGBColor(0xEF, 0xF6, 0xFF)),
    }
    text_col, _ = colors.get(color, colors['info'])
    p = doc.add_paragraph()
    label = p.add_run('EXAM TIP: ')
    label.bold = True
    label.font.color.rgb = text_col
    label.font.size = Pt(9.5)
    content = p.add_run(tip_text)
    content.font.color.rgb = text_col
    content.font.size = Pt(9.5)
    p.paragraph_format.left_indent = Cm(0.4)
    p.paragraph_format.space_before = Pt(2)
    p.paragraph_format.space_after = Pt(5)

def clean(text):
    """Remove control characters and NULL bytes that break lxml."""
    import re
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text).strip()

def add_slide_content(doc, slide, slide_num):
    """Extract and format a single slide into the document."""
    shapes = slide.shapes
    title_text = ''
    body_paras = []

    for shape in shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        paras = [clean(p.text) for p in tf.paragraphs if clean(p.text)]
        if not paras:
            continue

        # Heuristic: first non-empty text shape with large font = title
        if not title_text:
            title_text = clean(paras[0])
            rest = paras[1:]
        else:
            rest = paras

        for para_text in rest:
            body_paras.append(clean(para_text))

    if not title_text:
        return

    # Slide heading
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after = Pt(2)
    run = p.add_run(f'Slide {slide_num}  |  {title_text}')
    run.bold = True
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

    # Body bullets
    for text in body_paras:
        if not text:
            continue
        text = clean(text)
        if not text:
            continue
        # Detect code (contains indentation, def, import, #, =, etc.)
        is_code = any(tok in text for tok in ['import ', 'def ', '()', '  ', '>>>', '.filter(', '.map(', 'RDD', 'sc.', 'spark.', 'da.', 'dask', '.compute(', 'lambda', '    '])
        if is_code and len(text) > 10:
            p2 = doc.add_paragraph()
            r = p2.add_run(text)
            r.font.name = 'Courier New'
            r.font.size = Pt(8)
            r.font.color.rgb = RGBColor(0x7C, 0x3A, 0xED)
            p2.paragraph_format.left_indent = Cm(0.6)
            p2.paragraph_format.space_after = Pt(1)
        else:
            p2 = doc.add_paragraph(style='List Bullet')
            r = p2.add_run(text)
            r.font.size = Pt(10)
            p2.paragraph_format.left_indent = Cm(0.3)
            p2.paragraph_format.space_after = Pt(1)

    # Check for matching tip
    for key, tip in TIPS.items():
        if key.lower() in title_text.lower():
            add_tip(doc, tip, color='warn' if 'WARNING' in tip or 'EXAM' in tip else 'info')
            break

# ── Build document ────────────────────────────────────────────────────────────
prs = Presentation(str(PPTX))
slides = list(prs.slides)

doc = Document()
for section in doc.sections:
    section.top_margin    = Cm(2.0)
    section.bottom_margin = Cm(2.0)
    section.left_margin   = Cm(2.2)
    section.right_margin  = Cm(2.2)

# ── Cover ──
title_p = doc.add_heading('DSAI4205 — Exam Focus Notes', level=0)
title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER

sub = doc.add_paragraph('Dr. Fong\'s Exam Review Slides (146 slides) + Lecture Transcript Tips')
sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
sub.runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
sub.runs[0].font.size = Pt(11)

note = doc.add_paragraph()
r1 = note.add_run('Orange "EXAM TIP" boxes = directly sourced from Dr. Fong\'s Apr 11 lecture transcript')
r1.font.color.rgb = RGBColor(0x92, 0x40, 0x0E)
r1.italic = True
r1.font.size = Pt(10)
note.alignment = WD_ALIGN_PARAGRAPH.CENTER
doc.add_paragraph()

# ── Exam structure summary ──
doc.add_heading('Exam Structure (Dr. Fong confirmed)', level=1)
tbl = doc.add_table(rows=5, cols=2)
tbl.style = 'Table Grid'
rows_data = [
    ('Component', 'Details'),
    ('Section A', '22 MCQ, 1 mark each'),
    ('Section B', 'Long questions — scenario, computation, explanation'),
    ('Programming', '~10% only (RDD-based, NOT DataFrame). Skeleton with TO-DO gaps.'),
    ('Past paper', 'On BEPA — but Sem 1 had 30% programming + L10 in long Q. Coverage differs.'),
]
for ri, (a, b) in enumerate(rows_data):
    tbl.rows[ri].cells[0].text = a
    tbl.rows[ri].cells[1].text = b
    if ri == 0:
        set_cell_bg(tbl.rows[ri].cells[0], '1E3A5F')
        set_cell_bg(tbl.rows[ri].cells[1], '1E3A5F')
        for cell in tbl.rows[ri].cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    run.bold = True
                    run.font.size = Pt(10)
    else:
        for cell in tbl.rows[ri].cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10)
doc.add_paragraph()

doc.add_heading('What\'s in Long Question vs MCQ Only', level=1)
scope_tbl = doc.add_table(rows=13, cols=4)
scope_tbl.style = 'Table Grid'
scope_data = [
    ('Lecture', 'Long Q?', 'MCQ?', 'Notes'),
    ('L1 — 6 V\'s', 'YES', 'YES', 'Scenario-based passage → identify V\'s + propose solutions'),
    ('L2 — HDFS', 'Possibly', 'YES', '3-way replication, fault tolerance'),
    ('L2 — Dask', 'NO', '1–2 Qs', 'map_overlap, lazy, chunk size trade-offs'),
    ('L3 — MapReduce', 'YES ★', 'YES', 'Full flow: mapper → shuffle → reducer. Know intermediate formats.'),
    ('L3 — RDD', 'YES (10% prog)', 'YES', 'RDD only. Revise Tutorial 3. Skeleton code with TO-DO gaps.'),
    ('L3/L4 — DataFrame', 'NO', '3–4 Qs', 'Catalyst optimizer, key ops, vs RDD differences'),
    ('L5 — NLP', 'YES', 'YES', 'BPE (vocab=10), Levenshtein DP table, stemming vs lemmatization'),
    ('L6 — Embeddings', 'Conceptual', 'YES', 'NO equations. Skip-gram vs CBOW concepts, why softmax, M definition'),
    ('L7 — PageRank', 'YES ★', 'YES', 'Build matrix, power iteration, dead end + spider trap fixes'),
    ('L8 — Graph', 'YES ★', 'YES', 'Density, CC, betweenness (vertex only). Edge betweenness NOT tested.'),
    ('L9 — NoSQL/CAP', 'NO', 'YES', 'CAP theorem, BASE, strong vs eventual consistency'),
    ('L10 — RecSys', 'Not this sem', 'Possibly', 'Was in long Q in Sem 1. Not confirmed this semester.'),
]
for ri, row in enumerate(scope_data):
    for ci, val in enumerate(row):
        cell = scope_tbl.rows[ri].cells[ci]
        cell.text = val
        for para in cell.paragraphs:
            for run in para.runs:
                run.font.size = Pt(9)
                if ri == 0:
                    run.bold = True
        if ri == 0:
            set_cell_bg(cell, '1E3A5F')
            for para in cell.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        elif '★' in row[1]:
            set_cell_bg(cell, 'FEF3C7')
doc.add_paragraph()

doc.add_heading('Common Mistakes — Dr. Fong\'s Direct Warnings', level=1)
warnings = [
    ('Levenshtein DP table',      'Source word at TOP (columns). Target word at LEFT (rows). NEVER mix up. Wrong orientation = wrong entire question.'),
    ('PageRank adjacency matrix', 'Source node at TOP (columns). Destination node at LEFT (rows). Column sums must equal 1.'),
    ('CBOW window size M',        'M = half-window size. 2 words each side = M=2, NOT M=4. Divide by 2M in average.'),
    ('MapReduce shuffle',         'Must mention the shuffle/sort step between mapper and reducer. It is automatic but you must describe it in answers.'),
    ('Edge betweenness',          'NOT tested. Only vertex betweenness centrality is in scope. "Ash centrality will not be answered in your final exam."'),
]
warn_tbl = doc.add_table(rows=len(warnings)+1, cols=2)
warn_tbl.style = 'Table Grid'
warn_tbl.rows[0].cells[0].text = 'Topic'
warn_tbl.rows[0].cells[1].text = 'Warning'
set_cell_bg(warn_tbl.rows[0].cells[0], '7F1D1D')
set_cell_bg(warn_tbl.rows[0].cells[1], '7F1D1D')
for cell in warn_tbl.rows[0].cells:
    for para in cell.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.bold = True
            run.font.size = Pt(10)
for ri, (topic, warning) in enumerate(warnings, 1):
    warn_tbl.rows[ri].cells[0].text = topic
    warn_tbl.rows[ri].cells[1].text = warning
    set_cell_bg(warn_tbl.rows[ri].cells[0], 'FEE2E2')
    set_cell_bg(warn_tbl.rows[ri].cells[1], 'FFF7ED')
    for ci in range(2):
        for para in warn_tbl.rows[ri].cells[ci].paragraphs:
            for run in para.runs:
                run.font.size = Pt(9.5)

doc.add_page_break()

# ── Slide content by section ──
for section_title, slide_range in SECTIONS:
    doc.add_heading(section_title, level=1)
    for idx in slide_range:
        if idx - 1 < len(slides):
            add_slide_content(doc, slides[idx - 1], idx)
    doc.add_page_break()

doc.save(str(OUT))
print(f'Saved: {OUT}')

# Verify
from docx import Document as DDoc
d2 = DDoc(str(OUT))
non_empty = [p for p in d2.paragraphs if p.text.strip()]
print(f'Non-empty paragraphs: {len(non_empty)}')
print(f'Tables: {len(d2.tables)}')
